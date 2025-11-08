from django.shortcuts import render, redirect
from django.http import FileResponse
from PyPDF2 import PdfReader, PdfWriter
from io import BytesIO
import base64

def merge_pdfs(request):
    if "uploaded_pdfs" not in request.session:
        request.session["uploaded_pdfs"] = []
    uploaded_pdfs = request.session["uploaded_pdfs"]

    if request.method == "POST":
        # Upload PDF(s)
        if 'pdf_files' in request.FILES:
            for pdf in request.FILES.getlist("pdf_files"):
                file_bytes = pdf.read()
                # Store safely as Latin-1 for session (binary-safe)
                uploaded_pdfs.append({
                    "name": pdf.name,
                    "data": file_bytes.decode('latin1')
                })
            request.session["uploaded_pdfs"] = uploaded_pdfs
            request.session.modified = True
            return redirect("merge_pdfs")

        # Remove one file
        if 'remove' in request.POST:
            index = int(request.POST.get("remove"))
            if 0 <= index < len(uploaded_pdfs):
                uploaded_pdfs.pop(index)
                request.session["uploaded_pdfs"] = uploaded_pdfs
                request.session.modified = True
            return redirect("merge_pdfs")

        # Merge PDFs
        if 'merge' in request.POST:
            if len(uploaded_pdfs) < 2:
                return render(request, "myapp/merge_pdfs.html", {
                    "error": "Upload at least two PDF files before merging.",
                    "uploaded": uploaded_pdfs
                })

            merged = PdfWriter()
            for pdf_dict in uploaded_pdfs:
                pdf_stream = BytesIO(pdf_dict["data"].encode('latin1'))
                reader = PdfReader(pdf_stream)
                for page in reader.pages:
                    merged.add_page(page)

            merged_pdf = BytesIO()
            merged.write(merged_pdf)
            merged_pdf.seek(0)

            request.session["merged_pdf"] = merged_pdf.getvalue().decode('latin1')
            request.session["uploaded_pdfs"] = []
            request.session.modified = True

            return redirect("merged_download")

    # Convert to Base64 only when rendering preview
    preview_list = []
    for pdf in uploaded_pdfs:
        encoded = base64.b64encode(pdf["data"].encode('latin1')).decode('utf-8')
        preview_list.append({
            "name": pdf["name"],
            "data": encoded
        })

    return render(request, "myapp/merge_pdfs.html", {"uploaded": preview_list})




def merged_download(request):
    if "merged_pdf" not in request.session:
        return redirect("merge_pdfs")

    if request.method == "POST":
        # When user clicks "Download"
        pdf_data = request.session["merged_pdf"].encode('latin1')
        merged_pdf = BytesIO(pdf_data)
        del request.session["merged_pdf"]
        return FileResponse(merged_pdf, as_attachment=True, filename="merged.pdf")

    return render(request, "myapp/merged_download.html")


def home(request):
    return render(request, "myapp/home.html")


# split pdf views

from django.shortcuts import render, redirect
from django.http import FileResponse
from io import BytesIO
from PyPDF2 import PdfReader, PdfWriter
import zipfile
import re

def split_pdf(request):
    if request.method == "POST":
        pdf_file = request.FILES.get("pdf_file")
        ranges_text = request.POST.get("ranges")

        if not pdf_file or not ranges_text:
            return render(request, "myapp/split_pdf.html", {
                "error": "Please upload a PDF and enter valid page ranges."
            })

        # Parse page ranges safely
        try:
            parts = [p.strip() for p in ranges_text.split(',') if p.strip()]
            ranges = []
            for p in parts:
                if '-' in p:
                    start, end = map(int, p.split('-'))
                    ranges.append((start, end))
                else:
                    n = int(p)
                    ranges.append((n, n))
        except ValueError:
            return render(request, "myapp/split_pdf.html", {
                "error": "Invalid range format. Use like 1-3,5,7-9."
            })

        # Read the PDF
        try:
            reader = PdfReader(pdf_file)
        except:
            return render(request, "myapp/split_pdf.html", {
                "error": "The uploaded file is not a valid PDF."
            })

        total_pages = len(reader.pages)
        for (s, e) in ranges:
            if s < 1 or e > total_pages or s > e:
                return render(request, "myapp/split_pdf.html", {
                    "error": f"Invalid range {s}-{e}. PDF has {total_pages} pages."
                })

        # Create ZIP in memory
        zip_buffer = BytesIO()
        with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zf:
            for i, (start, end) in enumerate(ranges, start=1):
                writer = PdfWriter()
                for p in range(start - 1, end):
                    writer.add_page(reader.pages[p])
                part_buffer = BytesIO()
                writer.write(part_buffer)
                part_buffer.seek(0)
                zf.writestr(f"part_{i}_{start}-{end}.pdf", part_buffer.read())
        zip_buffer.seek(0)

        # Store split result in session temporarily
        request.session["split_zip"] = zip_buffer.getvalue().decode('latin1')

        return redirect("split_result")

    return render(request, "myapp/split_pdf.html")


def split_result(request):
    return render(request, "myapp/split_result.html")


def download_split(request):
    if "split_zip" not in request.session:
        return redirect("split_pdf")

    zip_data = request.session["split_zip"].encode('latin1')
    zip_buffer = BytesIO(zip_data)
    del request.session["split_zip"]
    return FileResponse(zip_buffer, as_attachment=True, filename="split_pdfs.zip")


# word to pdf views

from django.shortcuts import render
from django.http import JsonResponse, HttpResponse
from django.urls import reverse
from docx2pdf import convert
import tempfile
import os

def word_to_pdf(request):
    """Convert uploaded Word file to PDF (temporary, not saved in DB or media)."""
    if request.method == "GET":
        return render(request, "myapp/word_to_pdf.html")

    word_file = request.FILES.get("word")
    if not word_file:
        return JsonResponse({"error": "No file uploaded"}, status=400)

    # Save temporarily
    with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmp_word:
        tmp_word.write(word_file.read())
        tmp_word_path = tmp_word.name

    # Create a temporary output folder
    output_dir = tempfile.mkdtemp()

    try:
        # Convert Word to PDF (keeps original alignment & structure)
        convert(tmp_word_path, output_dir)

        # Get output file path
        pdf_path = os.path.join(
            output_dir,
            os.path.basename(tmp_word_path).replace(".docx", ".pdf")
        )

        # Read the PDF data into memory
        with open(pdf_path, "rb") as f:
            pdf_data = f.read()

        # Store temporarily in session
        request.session["pdf_data"] = pdf_data.decode("latin1")
        request.session["pdf_name"] = word_file.name.replace(".docx", ".pdf")

        # Redirect to download page
        redirect_url = reverse("word_to_pdf_download")
        return JsonResponse({"redirect_url": redirect_url})

    except Exception as e:
        return JsonResponse({"error": f"Conversion failed: {str(e)}"}, status=500)

    finally:
        # Cleanup temp files
        if os.path.exists(tmp_word_path):
            os.remove(tmp_word_path)
        if os.path.exists(output_dir):
            for f in os.listdir(output_dir):
                os.remove(os.path.join(output_dir, f))
            os.rmdir(output_dir)


def word_to_pdf_download(request):
    """Show download page after conversion."""
    file_name = request.session.get("pdf_name", "converted.pdf")
    pdf_data = request.session.get("pdf_data")

    if not pdf_data:
        return HttpResponse("<h3 style='text-align:center;color:red;'>No converted file found.</h3>")

    return render(request, "myapp/word_to_pdf_download.html", {"file_name": file_name})


def download_converted_pdf(request):
    """Send the converted PDF file for download."""
    pdf_data = request.session.get("pdf_data")
    file_name = request.session.get("pdf_name", "converted.pdf")

    if not pdf_data:
        return HttpResponse("<h3>No file available to download.</h3>")

    # Convert stored string data back to bytes
    pdf_bytes = pdf_data.encode("latin1")

    # Create downloadable response
    response = HttpResponse(pdf_bytes, content_type="application/pdf")
    response["Content-Disposition"] = f'attachment; filename="{file_name}"'

    return response



# compress pdf views

import io
import base64
from django.http import JsonResponse
from django.shortcuts import render
from PyPDF2 import PdfReader, PdfWriter


def compress_pdf_page(request):
    """Main compress page"""
    return render(request, "myapp/compress_pdf.html")


def compress_upload(request):
    """Handles PDF compression in memory"""
    if request.method == 'POST' and request.FILES.get('pdf'):
        pdf_file = request.FILES['pdf']
        quality = request.POST.get('quality', 'medium')

        try:
            # Read and write PDF
            reader = PdfReader(pdf_file)
            writer = PdfWriter()

            for page in reader.pages:
                writer.add_page(page)

            # Remove metadata
            writer.add_metadata({})

            # Write to memory
            output_buffer = io.BytesIO()
            writer.write(output_buffer)
            output_buffer.seek(0)
            pdf_bytes = output_buffer.read()

            # Convert to Base64 string
            pdf_base64 = base64.b64encode(pdf_bytes).decode('utf-8')

            # ✅ Save to session temporarily
            request.session['compressed_pdf'] = pdf_base64
            request.session['compressed_filename'] = f"compressed_{pdf_file.name}"

            return JsonResponse({'success': True})

        except Exception as e:
            return JsonResponse({'success': False, 'error': str(e)})

    return JsonResponse({'success': False, 'error': 'Invalid request'})


def compress_result_page(request):
    """Result / download page"""
    pdf_base64 = request.session.get('compressed_pdf')
    filename = request.session.get('compressed_filename', 'compressed.pdf')

    if not pdf_base64:
        return render(request, "myapp/compress_download.html", {
            'error': "No PDF found. Please compress a file first."
        })

    pdf_url = f"data:application/pdf;base64,{pdf_base64}"
    return render(request, "myapp/compress_download.html", {
        'filename': filename,
        'file_url': pdf_url
    })


# pdf to word views

from django.http import JsonResponse
from django.shortcuts import render
from pdf2docx import Converter
import tempfile, os, base64

def pdf_to_word_view(request):
    return render(request, "myapp/pdf_to_word.html")

def convert_upload(request):
    if request.method == "POST":
        pdf_file = request.FILES.get("pdf")
        if not pdf_file:
            return JsonResponse({"success": False, "error": "No file"})

        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_pdf:
            for chunk in pdf_file.chunks():
                temp_pdf.write(chunk)
            temp_pdf_path = temp_pdf.name

        word_path = temp_pdf_path.replace(".pdf", ".docx")

        try:
            cv = Converter(temp_pdf_path)
            cv.convert(word_path, start=0, end=None)
            cv.close()

            with open(word_path, "rb") as f:
                word_bytes = f.read()
            file_b64 = base64.b64encode(word_bytes).decode("utf-8")

            filename = os.path.basename(word_path)
            file_url = f"data:application/vnd.openxmlformats-officedocument.wordprocessingml.document;base64,{file_b64}"

            os.remove(temp_pdf_path)
            os.remove(word_path)

            return JsonResponse({"success": True, "filename": filename, "file_url": file_url})
        except Exception as e:
            return JsonResponse({"success": False, "error": str(e)})

    return JsonResponse({"success": False, "error": "Invalid request"})

def convert_result(request):
    filename = request.GET.get("file", "")
    file_url = request.GET.get("url", "")
    return render(request, "myapp/pdf_to_word_result.html", {"filename": filename, "file_url": file_url})


# pdf to ppt views

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
from django.http import JsonResponse
from django.shortcuts import render
import io
import base64

def pdf_to_ppt_page(request):
    return render(request, "myapp/pdf_to_ppt.html")

def pdf_to_ppt_convert(request):
    if request.method == "POST" and request.FILES.get("pdf"):
        try:
            pdf_file = request.FILES["pdf"]
            pdf_doc = fitz.open(stream=pdf_file.read(), filetype="pdf")

            prs = Presentation()
            blank_slide_layout = prs.slide_layouts[6]  # Empty slide

            for page in pdf_doc:
                # Convert each PDF page to image
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                img_data = io.BytesIO(pix.tobytes("png"))

                slide = prs.slides.add_slide(blank_slide_layout)
                slide.shapes.add_picture(img_data, Inches(0), Inches(0), Inches(10), Inches(7.5))

            output = io.BytesIO()
            prs.save(output)
            ppt_bytes = output.getvalue()
            output.close()

            filename = pdf_file.name.replace(".pdf", "_converted.pptx")

            # Encode file to Base64 to send via JSON (not saved to disk)
            ppt_b64 = base64.b64encode(ppt_bytes).decode("utf-8")

            return JsonResponse({
                "success": True,
                "filename": filename,
                "file_b64": ppt_b64
            })
        except Exception as e:
            print("Error:", e)
            return JsonResponse({"success": False, "error": str(e)})
    return JsonResponse({"success": False})


# image to pdf views

import io
from django.http import JsonResponse, FileResponse
from django.shortcuts import render
from PIL import Image
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4
from reportlab.lib.utils import ImageReader


def image_to_pdf_view(request):
    return render(request, "myapp/image_to_pdf.html")


def convert_image_to_pdf(request):
    if request.method == "POST":
        try:
            # Get multiple uploaded files
            images = request.FILES.getlist("images")
            if not images:
                return JsonResponse({"success": False, "error": "No images uploaded!"})

            # Create in-memory PDF
            pdf_buffer = io.BytesIO()
            c = canvas.Canvas(pdf_buffer, pagesize=A4)
            page_width, page_height = A4

            for img_file in images:
                img = Image.open(img_file).convert("RGB")
                temp = io.BytesIO()
                img.save(temp, format="PNG")
                temp.seek(0)
                image_reader = ImageReader(temp)

                img_width, img_height = img.size
                scale = min(page_width / img_width, page_height / img_height)
                new_width = img_width * scale
                new_height = img_height * scale
                x = (page_width - new_width) / 2
                y = (page_height - new_height) / 2

                c.drawImage(image_reader, x, y, new_width, new_height)
                c.showPage()

            c.save()
            pdf_buffer.seek(0)

            # Convert PDF to base64 for JS to download directly
            import base64
            pdf_b64 = base64.b64encode(pdf_buffer.getvalue()).decode("utf-8")

            return JsonResponse({"success": True, "pdf_b64": pdf_b64, "filename": "converted.pdf"})

        except Exception as e:
            print("❌ Conversion error:", e)
            return JsonResponse({"success": False, "error": "Error converting image!"})
    else:
        return JsonResponse({"success": False, "error": "Invalid request"})
  

#   pdf to image views

import base64
import io
import zipfile
import fitz  # PyMuPDF
from django.http import JsonResponse, HttpResponse
from django.shortcuts import render
from django.views.decorators.csrf import csrf_exempt

def pdf_to_image_page(request):
    return render(request, "myapp/pdf_to_image.html")

@csrf_exempt
def convert_pdf_to_image(request):
    try:
        if request.method != "POST":
            return JsonResponse({'success': False, 'error': 'Invalid request method'})

        pdf_files = request.FILES.getlist('pdfs')
        output_format = request.POST.get('format', 'png').lower()
        all_images = []

        for pdf_file in pdf_files:
            pdf_data = pdf_file.read()
            pdf_doc = fitz.open(stream=pdf_data, filetype="pdf")

            image_set = []
            for page in pdf_doc:
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # high quality
                img_data = pix.tobytes(output_format)
                img_b64 = base64.b64encode(img_data).decode('utf-8')
                image_set.append(f"data:image/{output_format};base64,{img_b64}")

            all_images.append({'filename': pdf_file.name, 'images': image_set})
            pdf_doc.close()

        return JsonResponse({'success': True, 'results': all_images, 'format': output_format})

    except Exception as e:
        print("💥 PDF to Image Conversion Error:", e)
        return JsonResponse({'success': False, 'error': str(e)})

@csrf_exempt
def download_zip(request):
    try:
        import json
        data = json.loads(request.body)
        images = data.get('images', [])
        output_format = data.get('format', 'png').lower()  # ✅ Added to handle JPG too

        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, 'w') as zip_file:
            for i, img_b64 in enumerate(images):
                header, img_data = img_b64.split(',', 1)
                img_bytes = base64.b64decode(img_data)
                zip_file.writestr(f'image_{i+1}.{output_format}', img_bytes)  # ✅ fixed extension

        zip_buffer.seek(0)
        response = HttpResponse(zip_buffer, content_type='application/zip')
        response['Content-Disposition'] = 'attachment; filename="converted_images.zip"'
        return response

    except Exception as e:
        print("💥 ZIP Error:", e)
        return JsonResponse({'success': False, 'error': str(e)})



# image compress views

from django.shortcuts import render
from django.http import JsonResponse
from django.views.decorators.csrf import csrf_exempt
from PIL import Image
import io, base64

def image_compress_page(request):
    return render(request, "myapp/image_compress.html")

@csrf_exempt
def compress_images(request):
    try:
        if request.method != "POST":
            return JsonResponse({'success': False, 'error': 'Invalid request method'})

        image_files = request.FILES.getlist('images')
        results = []

        for img_file in image_files:
            img = Image.open(img_file)
            buffer = io.BytesIO()
            img.save(buffer, format='JPEG', quality=60, optimize=True)
            compressed_data = buffer.getvalue()
            encoded_img = base64.b64encode(compressed_data).decode('utf-8')
            results.append({
                'filename': f"compressed_{img_file.name}",
                'image': f"data:image/jpeg;base64,{encoded_img}"
            })

        return JsonResponse({'success': True, 'results': results})

    except Exception as e:
        print("💥 Image Compression Error:", e)
        return JsonResponse({'success': False, 'error': str(e)})


# pdf protect views

import io
from django.http import JsonResponse
from django.shortcuts import render
from PyPDF2 import PdfReader, PdfWriter
import base64

def pdf_protect_page(request):
    return render(request, "myapp/pdf_protect.html")

def pdf_protect_convert(request):
    try:
        pdf_file = request.FILES.get("pdf")
        password = request.POST.get("password")

        if not pdf_file or not password:
            return JsonResponse({"success": False, "error": "Missing file or password"})

        # Read PDF directly from memory
        reader = PdfReader(pdf_file)
        writer = PdfWriter()

        for page in reader.pages:
            writer.add_page(page)

        writer.encrypt(password)

        # Write protected PDF to memory (no save on disk)
        protected_buffer = io.BytesIO()
        writer.write(protected_buffer)
        protected_buffer.seek(0)

        # Encode to Base64 for direct download
        protected_data = base64.b64encode(protected_buffer.getvalue()).decode("utf-8")

        return JsonResponse({
            "success": True,
            "filename": f"protected_{pdf_file.name}",
            "file_data": protected_data,
        })

    except Exception as e:
        print("Error:", e)
        return JsonResponse({"success": False, "error": str(e)})


# unlock pdf views

import io
import base64
from django.shortcuts import render
from django.http import JsonResponse
from PyPDF2 import PdfReader, PdfWriter

def pdf_unlock_page(request):
    return render(request, "myapp/pdf_unlock.html")

def pdf_unlock_convert(request):
    try:
        if request.method != "POST":
            return JsonResponse({"success": False, "error": "Invalid request method"})

        pdf_file = request.FILES.get("pdf")
        password = request.POST.get("password", "")

        if not pdf_file:
            return JsonResponse({"success": False, "error": "No file uploaded"})

        reader = PdfReader(pdf_file)

        if reader.is_encrypted:
            # Try to decrypt with provided password
            try:
                # PyPDF2: decrypt returns 0/1 or raises depending on version; handle both
                decrypt_ok = False
                try:
                    res = reader.decrypt(password)
                    # some versions return 1 for success
                    if res is True or res == 1 or res == 2:
                        decrypt_ok = True
                except Exception:
                    # fallback: try calling with password param when creating reader is not possible here
                    pass

                # if not decrypted yet, attempt another approach (some PyPDF2 versions)
                if not decrypt_ok:
                    # Re-open reader from bytes and pass password (works in some versions)
                    raw = pdf_file.read()
                    reader = PdfReader(io.BytesIO(raw))
                    if reader.is_encrypted:
                        try:
                            reader.decrypt(password)
                        except Exception:
                            pass

                if reader.is_encrypted:
                    # final check
                    try:
                        # Try to read first page to confirm
                        _ = reader.pages[0]
                    except Exception:
                        return JsonResponse({"success": False, "error": "Incorrect password or encrypted PDF."})
            except Exception:
                return JsonResponse({"success": False, "error": "Failed to decrypt. Wrong password?"})
        # If not encrypted, continue (we simply repackage without error)

        writer = PdfWriter()
        for p in reader.pages:
            writer.add_page(p)

        # Write unlocked PDF into memory
        out_buf = io.BytesIO()
        writer.write(out_buf)
        out_buf.seek(0)
        b64 = base64.b64encode(out_buf.read()).decode("utf-8")

        return JsonResponse({
            "success": True,
            "filename": f"unlocked_{pdf_file.name}",
            "file_data": b64
        })
    except Exception as e:
        print("Unlock Error:", e)
        return JsonResponse({"success": False, "error": str(e)})


# water mark views

from django.shortcuts import render
from django.http import FileResponse, JsonResponse
from PyPDF2 import PdfReader, PdfWriter
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from io import BytesIO

def pdf_watermark(request):
    return render(request, "myapp/pdf_watermark.html")

def preview_watermark(request):
    if request.method == "POST" and request.FILES.get("pdf"):
        pdf_file = request.FILES["pdf"]
        text = request.POST.get("text", "")
        position = request.POST.get("position", "center")
        opacity = float(request.POST.get("opacity", 0.3))

        try:
            # Create watermark overlay
            watermark_stream = BytesIO()
            c = canvas.Canvas(watermark_stream, pagesize=letter)
            c.setFont("Helvetica-Bold", 40)
            c.setFillGray(0, opacity)

            # ✅ Position logic
            if position == "center":
                c.drawCentredString(300, 400, text)
            elif position == "top left":
                c.drawString(50, 750, text)
            elif position == "bottom right":
                c.drawRightString(550, 50, text)
            elif position == "top diagonal":
                c.saveState()
                c.translate(310, 300)
                c.rotate(45)
                c.drawCentredString(0, 0, text)
                c.restoreState()
            elif position == "bottom diagonal":
                c.saveState()
                c.translate(310, 300)
                c.rotate(-45)
                c.drawCentredString(0, 0, text)
                c.restoreState()

            c.save()

            watermark_stream.seek(0)
            watermark_pdf = PdfReader(watermark_stream)
            original_pdf = PdfReader(pdf_file)
            output = PdfWriter()

            watermark_page = watermark_pdf.pages[0]

            for page in original_pdf.pages:
                page.merge_page(watermark_page)
                output.add_page(page)

            output_stream = BytesIO()
            output.write(output_stream)
            output_stream.seek(0)

            return FileResponse(
                output_stream,
                as_attachment=False,
                filename="preview.pdf",
                content_type="application/pdf",
            )
        except Exception as e:
            print("Error:", e)
            return JsonResponse({"error": str(e)}, status=500)
    return JsonResponse({"error": "Invalid request"}, status=400)


# sign pdf views

# myapp/views.py
import io
import fitz  # PyMuPDF
from PIL import Image
from django.http import HttpResponse
from django.shortcuts import render


def sign_pdf_page(request):
    return render(request, "myapp/sign_pdf.html")


def sign_pdf(request):
    try:
        pdf_file = request.FILES.get("pdf")
        if not pdf_file:
            return HttpResponse("No file uploaded", status=400)

        sig_type = request.POST.get("sigType")
        x = float(request.POST.get("x", 0))
        y = float(request.POST.get("y", 0))
        page_num = int(request.POST.get("page", 0))
        font_family = request.POST.get("fontFamily", "helv")  # safe font

        pdf_bytes = pdf_file.read()
        pdf_doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        if page_num >= len(pdf_doc):
            page_num = 0
        page = pdf_doc[page_num]

        # make sure coordinates inside page
        page_width, page_height = page.rect.width, page.rect.height
        x = max(0, min(x, page_width - 150))
        y = max(0, min(y, page_height - 50))

        if sig_type == "text":
            text = request.POST.get("sigText", "")
            font_size = int(request.POST.get("fontSize", 14))

            fonts = {
                "Times-Roman": "Times-Roman",
                "Helvetica": "helv",
                "Courier": "cour",
                "ZapfDingbats": "zapfdingbats",
                "Symbol": "symb"
            }
            chosen_font = fonts.get(font_family, "helv")

            page.insert_text(
                (x, y),
                text,
                fontsize=font_size,
                fontname=chosen_font,
                color=(0, 0, 0)
            )

        elif sig_type == "image":
            sig_img = request.FILES.get("sigImage")
            if sig_img:
                image = Image.open(sig_img).convert("RGB")
                img_byte = io.BytesIO()
                image.save(img_byte, format="PNG")
                rect = fitz.Rect(x, y, x + 150, y + 50)
                page.insert_image(rect, stream=img_byte.getvalue())

        output = io.BytesIO()
        pdf_doc.save(output)
        pdf_doc.close()
        output.seek(0)

        response = HttpResponse(output.getvalue(), content_type="application/pdf")
        response["Content-Disposition"] = 'attachment; filename="signed.pdf"'
        return response

    except Exception as e:
        print("❌ Error:", e)
        return HttpResponse("Conversion failed!", status=500) 



# pdf to excel views

# views.py (add these imports at top if not present)
import io
import pdfplumber
import pytesseract
from pdf2image import convert_from_bytes
from openpyxl import Workbook
from openpyxl.utils import get_column_letter
from openpyxl.styles import Alignment
from django.http import HttpResponse, JsonResponse
from django.shortcuts import render

# ✅ No file path needed — works directly from uploaded memory file

def pdf_to_excel_page(request):
    return render(request, "myapp/pdf_to_excel.html")


def convert_pdf_to_excel(request):
    if request.method != "POST":
        return JsonResponse({"success": False, "error": "Invalid request"}, status=400)

    try:
        pdf_file = request.FILES.get("pdf")
        if not pdf_file:
            return JsonResponse({"success": False, "error": "No file uploaded"}, status=400)

        # 📄 Read uploaded file directly from memory
        pdf_bytes = pdf_file.read()
        extracted_data = []

        # 1️⃣ Extract tables/text using pdfplumber
        with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
            for page in pdf.pages:
                table = page.extract_table()
                if table:
                    extracted_data.extend(table)
                else:
                    text = page.extract_text()
                    if text:
                        for line in text.split("\n"):
                            extracted_data.append([line])

        # 2️⃣ OCR fallback if no text extracted
        if not extracted_data:
            pages = convert_from_bytes(pdf_bytes)
            for page in pages:
                text = pytesseract.image_to_string(page)
                for line in text.split("\n"):
                    if line.strip():
                        extracted_data.append([line.strip()])

        if not extracted_data:
            return JsonResponse({"success": False, "error": "No readable text found"}, status=400)

        # 3️⃣ Create Excel
        wb = Workbook()
        ws = wb.active
        ws.title = "Extracted Data"

        for r, row in enumerate(extracted_data, start=1):
            for c, val in enumerate(row, start=1):
                ws.cell(row=r, column=c).value = val

        # 4️⃣ Auto adjust columns + wrap text
        for col in ws.columns:
            max_length = 0
            col_letter = get_column_letter(col[0].column)
            for cell in col:
                cell.alignment = Alignment(wrap_text=True, vertical="top")
                if cell.value:
                    max_length = max(max_length, len(str(cell.value)))
            ws.column_dimensions[col_letter].width = min(max_length + 2, 50)

        # 5️⃣ Save Excel to memory and return as response
        output = io.BytesIO()
        wb.save(output)
        output.seek(0)

        response = HttpResponse(
            output.getvalue(),
            content_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        )
        response["Content-Disposition"] = f'attachment; filename="{pdf_file.name.replace(".pdf", ".xlsx")}"'
        return response

    except Exception as e:
        print("❌ Conversion Error:", str(e))
        return JsonResponse({"success": False, "error": str(e)}, status=400)






# excel to pdf views

import openpyxl
from io import BytesIO
from django.shortcuts import render
from django.http import JsonResponse, HttpResponse
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, PageBreak
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet


def excel_to_pdf_view(request):
    if request.method == "GET":
        return render(request, "myapp/excel_to_pdf.html")

    excel_file = request.FILES.get("excel")
    if not excel_file:
        return JsonResponse({"error": "No file uploaded"}, status=400)

    try:
        # Load Excel workbook
        wb = openpyxl.load_workbook(excel_file, data_only=True)
        out = BytesIO()
        doc = SimpleDocTemplate(
            out,
            pagesize=landscape(A4),
            leftMargin=20,
            rightMargin=20,
            topMargin=20,
            bottomMargin=20
        )
        styles = getSampleStyleSheet()
        story = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            story.append(Paragraph(f"Sheet: {sheet_name}", styles["Heading3"]))
            story.append(Spacer(1, 8))

            # Read Excel data
            data = []
            for row in ws.iter_rows(values_only=True):
                data.append([str(cell) if cell is not None else "" for cell in row])

            if not data:
                story.append(Paragraph("(Empty sheet)", styles["Normal"]))
                continue

            # Calculate Excel-like widths
            col_widths = []
            for col in ws.columns:
                max_len = max((len(str(cell.value)) if cell.value else 1 for cell in col))
                col_widths.append(max_len * 6)

            # Auto-fit columns within page
            total_width = sum(col_widths)
            available_width = landscape(A4)[0] - 40
            scale = available_width / total_width if total_width > available_width else 1.0
            col_widths = [w * scale for w in col_widths]

            # Split large sheet into multiple pages
            rows_per_page = 40
            for i in range(0, len(data), rows_per_page):
                table_chunk = data[i:i + rows_per_page]
                table = Table(table_chunk, colWidths=col_widths, repeatRows=1)
                style = TableStyle([
                    ("GRID", (0, 0), (-1, -1), 0.25, colors.grey),
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#e8e8e8")),
                    ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                    ("FONTSIZE", (0, 0), (-1, -1), 8),
                    ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                    ("ALIGN", (0, 0), (-1, -1), "CENTER"),
                ])
                table.setStyle(style)
                story.append(table)
                story.append(PageBreak())

        # Build the PDF
        doc.build(story)
        out.seek(0)
        pdf_bytes = out.getvalue()
        pdf_name = excel_file.name.rsplit(".", 1)[0] + ".pdf"

        # Store session for download
        request.session["excel_to_pdf_bytes"] = pdf_bytes.hex()
        request.session["excel_to_pdf_name"] = pdf_name

        return JsonResponse({
            "pdf_name": pdf_name,
            "download_url": "/excel-to-pdf/file/"
        })

    except Exception as e:
        print("Error:", e)
        return JsonResponse({"error": str(e)}, status=500)


def excel_to_pdf_file(request):
    pdf_hex = request.session.get("excel_to_pdf_bytes")
    pdf_name = request.session.get("excel_to_pdf_name", "converted.pdf")

    if not pdf_hex:
        return HttpResponse("No file found", status=404)

    pdf_bytes = bytes.fromhex(pdf_hex)
    response = HttpResponse(pdf_bytes, content_type="application/pdf")
    response["Content-Disposition"] = f'attachment; filename="{pdf_name}"'
    return response









# ppt to pdf views

import os
import tempfile
import subprocess
from django.http import HttpResponse, JsonResponse
from django.shortcuts import render
import platform


def ppt_to_pdf_page(request):
    return render(request, "myapp/ppt_to_pdf.html")


def convert_ppt_to_pdf(request):
    if request.method != "POST":
        return JsonResponse({"success": False, "error": "Invalid request method"}, status=405)

    ppt_file = request.FILES.get("ppt")
    if not ppt_file:
        return JsonResponse({"success": False, "error": "No file uploaded"}, status=400)

    temp_ppt = None
    output_pdf = None

    try:
        # Save uploaded file temporarily
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
            for chunk in ppt_file.chunks():
                tmp.write(chunk)
            temp_ppt = tmp.name

        system_os = platform.system().lower()

        # For Windows, use PowerPoint COM if available
        if "windows" in system_os:
            try:
                import comtypes.client
                powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
                powerpoint.Visible = 1
                pdf_path = os.path.splitext(temp_ppt)[0] + ".pdf"
                presentation = powerpoint.Presentations.Open(temp_ppt)
                presentation.SaveAs(pdf_path, 32)
                presentation.Close()
                powerpoint.Quit()
                output_pdf = pdf_path
            except Exception as e:
                print("⚠️ PowerPoint not found, trying LibreOffice:", e)
                output_dir = tempfile.gettempdir()
                subprocess.run(
                    ["libreoffice", "--headless", "--convert-to", "pdf", "--outdir", output_dir, temp_ppt],
                    check=True
                )
                output_pdf = os.path.splitext(temp_ppt)[0] + ".pdf"

        # For Linux/Mac use LibreOffice
        else:
            output_dir = tempfile.gettempdir()
            subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "pdf", "--outdir", output_dir, temp_ppt],
                check=True
            )
            output_pdf = os.path.splitext(temp_ppt)[0] + ".pdf"

        if not os.path.exists(output_pdf):
            raise Exception("Conversion failed: PDF not created.")

        with open(output_pdf, "rb") as f:
            pdf_data = f.read()

        response = HttpResponse(pdf_data, content_type="application/pdf")
        response["Content-Disposition"] = 'attachment; filename="converted_ppt.pdf"'
        return response

    except subprocess.CalledProcessError as e:
        print("❌ LibreOffice failed:", e)
        return JsonResponse({"success": False, "error": "LibreOffice conversion failed."}, status=500)

    except Exception as e:
        print("❌ Error:", e)
        return JsonResponse({"success": False, "error": str(e)}, status=500)

    finally:
        if temp_ppt and os.path.exists(temp_ppt):
            os.remove(temp_ppt)
        if output_pdf and os.path.exists(output_pdf):
            os.remove(output_pdf)



from django.shortcuts import render, redirect
from django.contrib.auth.models import User
from django.contrib import messages
from django.contrib.auth import login

def signup_page(request):
    if request.method == "POST":
        name = request.POST.get("name").strip()
        email = request.POST.get("email").strip()
        password = request.POST.get("password").strip()
        confirm_password = request.POST.get("confirm_password").strip()

        # 1️⃣ Password match check
        if password != confirm_password:
            messages.error(request, "Passwords do not match!")
            return redirect('signup_page')

        # 2️⃣ Email already exists check
        if User.objects.filter(email=email).exists():
            messages.error(request, "Email already registered. Please login.")
            return redirect('signup_page')

        # 3️⃣ Username already exists check
        if User.objects.filter(username=name).exists():
            messages.error(request, "Username already taken. Try another.")
            return redirect('signup_page')

        # 4️⃣ Create new user
        user = User.objects.create_user(username=name, email=email, password=password)
        user.save()

        # 5️⃣ Auto login & redirect
        login(request, user)
        messages.success(request, "Signup successful! Welcome to VetriFiles.")
        return redirect('home')

    return render(request, "myapp/signup.html")


from django.shortcuts import render, redirect
from django.contrib import messages
from django.contrib.auth import authenticate, login
from django.contrib.auth.models import User

def login_page(request):
    if request.method == "POST":
        username_or_email = request.POST.get("username_or_email").strip()
        password = request.POST.get("password").strip()

        # Allow login by username or email
        user = None
        if User.objects.filter(username=username_or_email).exists():
            user = authenticate(username=username_or_email, password=password)
        elif User.objects.filter(email=username_or_email).exists():
            user_obj = User.objects.get(email=username_or_email)
            user = authenticate(username=user_obj.username, password=password)

        if user is not None:
            login(request, user)
            messages.success(request, "Login successful! Welcome back.")
            return redirect('home')
        else:
            messages.error(request, "Invalid username/email or password.")

    return render(request, "myapp/login.html")


# 🔹 Password reset (manual simple version)
def reset_password(request):
    if request.method == "POST":
        password = request.POST.get("password").strip()
        confirm_password = request.POST.get("confirm_password").strip()
        email = request.POST.get("email").strip()

        if password != confirm_password:
            messages.error(request, "Passwords do not match!")
            return redirect('login_page')

        try:
            user = User.objects.get(email=email)
            user.set_password(password)
            user.save()
            messages.success(request, "Password updated successfully. Please login.")
            return redirect('login_page')
        except User.DoesNotExist:
            messages.error(request, "Email not found. Try again.")
            return redirect('login_page')

    return redirect('login_page')
